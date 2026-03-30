from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

from src.common.config import DATA_PROCESSED, DECK, REPORTS, SECTION3_OUTPUT_CHARTS

CHARTS = SECTION3_OUTPUT_CHARTS
from src.common.utils import write_markdown


SLIDES = [
    (
        "HDB Resale Buyer Decision Support Can Be Structured Around Three Questions",
        [
            "How has the resale market evolved overall and by town?",
            "Where can a buyer buy the most space for a given budget?",
            "Which towns balance affordability with accessibility best?",
        ],
    ),
    (
        "Prices And Transactions Show A Clear Long-Run Uptrend With Strong Town Dispersion",
        [
            "Use the market overview dashboard to compare national and town-level transaction counts and median prices.",
            "Keep flat type as a first-class filter because town price rankings shift materially by unit mix.",
        ],
    ),
    (
        "Budget-Constrained Buyers Can Trade Location For Space In A Measurable Way",
        [
            "The budget dashboard ranks towns by attainable floor area for common budget bands.",
            "This turns anecdotal affordability conversations into a concrete shortlist.",
        ],
    ),
    (
        "A Location Score Lets Buyers Balance Price, Size, MRT Access, And CBD Distance",
        [
            "Town-centroid proximity measures are enough for a strong interview prototype and can later be upgraded to block-level geocoding.",
            "The location dashboard is designed to surface tradeoffs rather than chase a single 'best' town.",
        ],
    ),
    (
        "A Tight 2014 Modeling Benchmark Gives A Defensible Baseline For Market Expectations",
        [
            "Three candidate models are compared using only flat type, flat age, and town, as requested.",
            "The best family is then reused for case scoring and discussed in terms of both accuracy and explainability.",
        ],
    ),
    (
        "The Nov 2017 Yishun Example Should Be Framed As A Model-Based Reasonableness Check, Not A Verdict In Isolation",
        [
            "Present the predicted price, residual, and confidence band in plain language.",
            "Senior management should hear what the model knows, what it does not know, and how confident we are.",
        ],
    ),
    (
        "Policy Questions Need Careful Definitions, Controlled Comparisons, And Explicit Caveats",
        [
            "Yishun should be framed as a value town, not automatically the absolute cheapest town.",
            "Flat-size shrinkage appears real in larger flat types after 2008, while the DTL2 and COE questions need especially careful causal interpretation.",
        ],
    ),
    (
        "The Recommendation Is To Launch A Buyer-Facing Dashboard And Use The Analytics Stack Internally For Monitoring",
        [
            "Dashboard outputs can support buyer transparency immediately.",
            "The same pipeline can be extended for policy evaluation, anomaly detection, and richer transaction scoring.",
        ],
    ),
]


def build_slide_outline() -> Path:
    lines = [
        "# Slide Outline",
        "",
        "This outline is written for a one-hour senior-management presentation with embedded speaker notes.",
        "",
    ]
    for index, (title, notes) in enumerate(SLIDES, start=1):
        lines.append(f"## Slide {index}: {title}")
        lines.append("")
        for note in notes:
            lines.append(f"- {note}")
        lines.append("")
    lines.extend(
        [
            "## Backup Slides",
            "",
            "- Detailed model comparison table with MAE, RMSE, and R2.",
            "- Flat-type classification report by class.",
            "- Yishun value story: raw price, price per sqm, and floor-area comparison.",
            "- Flat-size evidence focused on 4-room, 5-room, and executive trends after 2008.",
            "- Downtown Line Stage 2: headline DiD result plus pre-trend caution slide.",
            "- COE question: visual co-movement versus controlled relative-effect slide.",
        ]
    )
    path = REPORTS / "slide_outline.md"
    write_markdown(path, lines)
    return path


def _add_text_slide(prs: Presentation, title: str, bullets: list[str], image_path: Path | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    title_frame = slide.shapes.title.text_frame
    title_frame.paragraphs[0].font.size = Pt(26)

    body = slide.placeholders[1].text_frame
    body.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(18)

    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(5.6), Inches(1.6), width=Inches(4.0))


def generate_powerpoint() -> Path:
    DECK.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "HDB Resale Strategy Case"
    title_slide.placeholders[1].text = "Data pipeline, buyer dashboards, modeling, and policy analysis"

    chart_lookup = [
        CHARTS / "yishun_cheapest.png",
        CHARTS / "flat_sizes_over_time.png",
        CHARTS / "dtl2_did.png",
        CHARTS / "coe_link.png",
    ]
    for idx, (title, bullets) in enumerate(SLIDES):
        image = chart_lookup[idx - 4] if 4 <= idx < 8 and (idx - 4) < len(chart_lookup) else None
        _add_text_slide(prs, title, bullets, image_path=image)

    output_path = DECK / "hdb_case_interview_deck.pptx"
    prs.save(output_path)
    return output_path


def build_deck_artifacts() -> dict[str, str]:
    dataset_path = DATA_PROCESSED / "hdb_resale_processed.parquet"
    if not dataset_path.exists():
        raise FileNotFoundError("Processed dataset missing. Run `python -m src.pipeline.build_resale_analysis_dataset` first.")

    _ = pd.read_parquet(dataset_path, columns=["transaction_year"])
    outline = build_slide_outline()
    pptx = generate_powerpoint()
    return {"slide_outline": str(outline), "pptx": str(pptx)}
